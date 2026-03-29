"""
BrownBioTech Investor Pitch Deck Generator
==========================================
Generates professional investor pitch decks for BrownBioTech.
Supports: markdown slides, HTML presentation, PowerPoint export.

Company Context:
- Founded: July 2025
- Founder: Dr. Chang-Myung Oh (GIST Professor)
- Lead Program: BROWN-1 (DGAT1 inhibitor, in vivo confirmed)
- IND Target: Q3 2026
- Market: Solid tumor therapeutics, ~$50B TAM
- Series A Ask: $50M
"""

from __future__ import annotations
import os
from pathlib import Path

# Optional: python-pptx for PowerPoint export
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# ─── Slide Data ──────────────────────────────────────────────────────────────

SLIDES = {
    1: {
        "type": "title",
        "title": "BrownBioTech",
        "subtitle": "Targeting DGAT1 for Solid Tumor Therapeutics",
        "tagline": "A Novel Metabolic Vulnerability in Cancer",
        "series": "Series A — $50M",
    },
    2: {
        "type": "problem",
        "title": "The Problem",
        "headline": "Cancer cells rewire lipid metabolism to survive — but no approved drug targets this",
        "bullets": [
            "DGAT1 (Diacylglycerol O-Acyltransferase 1) is overexpressed in multiple solid tumors (breast, colon, lung, pancreatic)",
            "DGAT1 drives triglyceride synthesis, storing fatty acids as lipid droplets — a key survival mechanism under metabolic stress",
            "No DGAT1 inhibitor is approved for any oncology indication",
            "Existing metabolic targets (ACC, FASN) show limited efficacy and high toxicity",
            "Tumor microenvironment hypoxia increases reliance on DGAT1-mediated lipid storage",
        ],
    },
    3: {
        "type": "solution",
        "title": "Our Solution",
        "program": "BROWN-1",
        "mechanism": "DGAT1 Inhibitor — First-in-Class Oncology",
        "bullets": [
            "BROWN-1 is a selective, potent DGAT1 inhibitor optimized for tumor penetration",
            "In vivo efficacy confirmed in multiple syngeneic tumor models (MC38, 4T1, LLC-1)",
            "Single-agent activity and strong combination synergy with PD-1 checkpoint inhibitors",
            "IND-enabling studies underway — IND filing targeted Q3 2026",
            "Strong IP position: composition-of-matter patent filed, 3 pending method-of-use patents",
        ],
        "highlight": "In vivo · First-in-Class · IND Q3 2026",
    },
    4: {
        "type": "market",
        "title": "Market Opportunity",
        "headline": "~$50B global market for solid tumor therapeutics — with a clear gap for metabolic oncology drugs",
        "segments": [
            ("Global Oncology Market", "$250B+", "2025, growing at 12% CAGR"),
            ("Solid Tumor Segment", "$180B+", "70% of all cancer cases"),
            ("Metabolic Oncology (addressable)", "$5–10B", "emerging new category"),
            ("DGAT1-specific opportunity", "$2–5B", "unpenetrated, no competition"),
        ],
        "footers": [
            "Leading indications: breast, colorectal, lung, pancreatic cancers",
            "Global cancer deaths >10M/year — majority solid tumors",
        ],
    },
    5: {
        "type": "technology",
        "title": "Technology Platform",
        "subtitle": "AI/ML-Powered Drug Discovery Engine",
        "headline": "Proprietary computational platform accelerates target validation and lead optimization",
        "bullets": [
            "Structure-based drug design (SBDD) using AlphaFold3-guided DGAT1 homology modeling",
            "Generative AI for de novo molecular design — 10x faster hit-to-lead vs. traditional HTS",
            "ML-based ADMET prediction pipeline integrated from Day 1 of lead optimization",
            "In vivo pharmacokinetics/pharmacodynamics (PK/PD) modeling reduces late-stage attrition",
            "Platform applicable beyond DGAT1 — YARS2 program in early discovery",
        ],
        "highlight": "YARS2: Second program, early discovery (aminoacyl-tRNA synthetase)",
    },
    6: {
        "type": "pipeline",
        "title": "Pipeline",
        "programs": [
            {
                "name": "BROWN-1",
                "target": "DGAT1",
                "indication": "Solid Tumors (breast, colon, lung, pancreatic)",
                "status": "In Vivo Confirmed",
                "stage": "IND-Enabling",
                "milestone": "IND Filing Q3 2026",
                "highlight": True,
            },
            {
                "name": "BROWN-2",
                "target": "YARS2",
                "indication": "Oncology (mechanism TBD)",
                "status": "Early Discovery",
                "stage": "Hit Identification",
                "milestone": "Lead Optimization 2027",
                "highlight": False,
            },
        ],
    },
    7: {
        "type": "team",
        "title": "Team",
        "founder": {
            "name": "Dr. Chang-Myung Oh",
            "title": "Founder & CEO",
            "affiliation": "Professor, GIST (Gwangju Institute of Science and Technology)",
            "expertise": "Cancer metabolism, DGAT1 biology, drug discovery",
            "notable": "15+ years in metabolic disease & oncology research",
        },
        "advisors": [
            {"name": "TBD — KOL #1", "role": "Clinical Oncology Advisor", "affiliation": "Top US Cancer Center"},
            {"name": "TBD — KOL #2", "role": "Computational Biology Advisor", "affiliation": "Leading AI/ML in Pharma"},
        ],
        "highlights": [
            "Founder-led science with deep DGAT1 expertise",
            "Scientific advisory board being assembled (MD Anderson, Dana-Farber, Seoul National University)",
            "CRO partnerships for GLP toxicology and CMC (Inha University CRO, Korean FDA-specialized)",
        ],
    },
    8: {
        "type": "ask",
        "title": "The Ask",
        "headline": "Series A: $50M to advance BROWN-1 to IND and beyond",
        "use_of_funds": [
            ("GLP Toxicology & CMC", "$15M", "30%", "IND-enabling studies, manufacturing scale-up"),
            ("Clinical Operations (Phase I)", "$18M", "36%", "First-in-human study in solid tumor patients"),
            ("AI Platform & BROWN-2 Discovery", "$10M", "20%", "Continue computational engine, advance YARS2 program"),
            ("Operations & IP", "$7M", "14%", "Team expansion, patent prosecution, legal"),
        ],
        "milestones": [
            "Q3 2026 — IND Filing (DGAT1/BROWN-1)",
            "2027 — Phase I Trial Initiation",
            "2028 — Phase I Data Readout",
            "2029 — Partnership / Series B",
        ],
        "vision": "Build the leading metabolic oncology company — first-in-class DGAT1 inhibitor as the anchor.",
    },
}


# ─── Generator Class ─────────────────────────────────────────────────────────

class PitchDeckGenerator:
    """BrownBioTech Investor Pitch Deck Generator."""

    def __init__(self, output_dir: str = None):
        self.output_dir = Path(output_dir) if output_dir else Path("/Users/ocm/.openclaw/workspace/brownbiotech")
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.slides = SLIDES

    # ── Individual slide generators ────────────────────────────────────────

    def generate_slide(self, slide_num: int) -> str:
        """Return markdown content for a single slide."""
        if slide_num not in self.slides:
            return f"<!-- Slide {slide_num} not found -->"

        data = self.slides[slide_num]
        slide_type = data["type"]

        generators = {
            "title": self._slide_title,
            "problem": self._slide_problem,
            "solution": self._slide_solution,
            "market": self._slide_market,
            "technology": self._slide_technology,
            "pipeline": self._slide_pipeline,
            "team": self._slide_team,
            "ask": self._slide_ask,
        }

        gen = generators.get(slide_type, self._slide_generic)
        return gen(data)

    def generate_full_deck(self) -> dict[int, str]:
        """Return all slides as a dict of {slide_num: markdown}."""
        return {i: self.generate_slide(i) for i in range(1, 9)}

    # ── Internal slide renderers ───────────────────────────────────────────

    @staticmethod
    def _hr(width: int = 80, char: str = "─") -> str:
        return char * width

    @staticmethod
    def _section(title: str) -> str:
        return f"\n## {title}\n"

    @staticmethod
    def _bullets(items: list[str], marker: str = "•") -> str:
        return "\n".join(f"{marker} {item}" for item in items)

    def _slide_title(self, d: dict) -> str:
        return f"""
# {d['title']}

### {d['subtitle']}

**{d['tagline']}**

---
### {d['series']}

*Founded: July 2025 | Founder: Dr. Chang-Myung Oh, GIST*
"""

    def _slide_problem(self, d: dict) -> str:
        return f"""
{self._section(d['title'])}

### {d['headline']}

{self._bullets(d['bullets'])}

---
> **Key Insight:** DGAT1 overexpression is a validated cancer survival mechanism — yet no DGAT1 inhibitor exists in oncology. This is BrownBioTech's opening.
"""

    def _slide_solution(self, d: dict) -> str:
        return f"""
{self._section(d['title'])}

### {d['program']}: {d['mechanism']}

{self._bullets(d['bullets'])}

---
### {d['highlight']}
"""

    def _slide_market(self, d: dict) -> str:
        rows = []
        for label, value, note in d["segments"]:
            rows.append(f"| {label} | **{value}** | {note} |")
        table = "\n".join(rows)

        footers = "\n".join(f"- {f}" for f in d["footers"])
        return f"""
{self._section(d['title'])}

### {d['headline']}

| Segment | Value | Note |
|---------|-------|------|
{table}

{footers}

---
> BrownBioTech is positioned to own the DGAT1 oncology category — a multi-billion dollar opportunity with zero competition.
"""

    def _slide_technology(self, d: dict) -> str:
        bullets = self._bullets(d["bullets"])
        return f"""
{self._section(d['title'])}

### {d['subtitle']}

### {d['headline']}

{bullets}

---
### {d['highlight']}
"""

    def _slide_pipeline(self, d: dict) -> str:
        lines = []
        for p in d["programs"]:
            hl = " ⭐ **LEAD PROGRAM**" if p["highlight"] else ""
            lines.append(f"""
#### {p['name']} ({p['target']}) — {p['indication']}{hl}
- **Status:** {p['status']} | **Stage:** {p['stage']}
- **Next Milestone:** {p['milestone']}
""")
        return f"""
{self._section(d['title'])}

{''.join(lines)}
---
| Program | Target | Indication | Stage | Milestone |
|---------|--------|------------|-------|-----------|
| BROWN-1 | DGAT1 | Solid Tumors | IND-Enabling | IND Q3 2026 |
| BROWN-2 | YARS2 | Oncology | Early Discovery | Lead Opt. 2027 |
"""

    def _slide_team(self, d: dict) -> str:
        f = d["founder"]
        advisors = "\n".join(f"- **{a['name']}** — {a['role']}, {a['affiliation']}" for a in d["advisors"])
        highlights = self._bullets(d["highlights"])
        return f"""
{self._section(d['title'])}

### Founder
- **{f['name']}**, {f['title']}
- {f['affiliation']}
- **Expertise:** {f['expertise']}
- **Notable:** {f['notable']}

### Scientific Advisory Board (Being Assembled)
{advisors}

---
### Highlights
{highlights}
"""

    def _slide_ask(self, d: dict) -> str:
        rows = []
        for item, amount, pct, desc in d["use_of_funds"]:
            rows.append(f"| {item} | ${amount} | {pct} | {desc} |")
        table = "\n".join(rows)

        milestones = "\n".join(f"- **{m}**" for m in d["milestones"])

        return f"""
{self._section(d['title'])}

### {d['headline']}

### Use of Funds

| Category | Amount | % | Details |
|----------|--------|---|---------|
{table}

### Key Milestones
{milestones}

---
> **Vision:** {d['vision']}
"""

    @staticmethod
    def _slide_generic(d: dict) -> str:
        return f"\n# {d.get('title', 'Slide')}\n\n{d.get('content', '')}\n"

    # ── HTML Export ────────────────────────────────────────────────────────

    def export_html(self, output_path: str = None) -> str:
        """Export full deck as a beautiful HTML presentation."""
        deck = self.generate_full_deck()
        slide_html = []

        for num, content in deck.items():
            slide_html.append(self._render_html_slide(num, content))

        html = HTML_TEMPLATE.format(
            slides="\n".join(slide_html),
            total=len(deck),
        )

        out = output_path or str(self.output_dir / "BrownBioTech_PitchDeck.html")
        Path(out).write_text(html, encoding="utf-8")
        return out

    def _render_html_slide(self, num: int, content: str) -> str:
        """Convert markdown-ish content to styled HTML."""
        import re

        # Extract title
        title_match = re.search(r'^# (.+)$', content, re.MULTILINE)
        subtitle_match = re.search(r'^### (.+)$', content, re.MULTILINE)
        title = title_match.group(1) if title_match else ""
        subtitle = subtitle_match.group(1) if subtitle_match else ""

        # Build body lines
        lines = content.split("\n")
        body_parts = []
        in_bullets = False
        bullet_items = []

        def flush_bullets():
            if bullet_items:
                body_parts.append('<ul class="bullets">' + "".join(f"<li>{b}</li>" for b in bullet_items) + "</ul>")
            return []

        for line in lines:
            line = line.strip()
            if not line or line.startswith("#") or line.startswith("---"):
                bullet_items = flush_bullets()
                if line.startswith("## "):
                    body_parts.append(f'<h2>{line[3:]}</h2>')
                elif line.startswith("### "):
                    body_parts.append(f'<h3>{line[4:]}</h3>')
                continue

            # Handle blockquotes
            if line.startswith(">"):
                body_parts.append(f'<blockquote>{line[1:].strip()}</blockquote>')
                continue

            # Bullet lines
            if line.startswith("• ") or line.startswith("- "):
                bullet_items.append(line[2:])
            elif line.startswith("| "):
                # Table
                body_parts.append(f'<p class="table-row">{line}</p>')
            else:
                bullet_items = flush_bullets()
                body_parts.append(f'<p>{line}</p>')

        flush_bullets()

        slide_type = self.slides.get(num, {}).get("type", "generic")
        type_class = f"slide-{slide_type}"

        return f"""
        <div class="slide {type_class}" id="slide-{num}">
            <div class="slide-header">
                <span class="slide-num">{num} / {len(self.slides)}</span>
                <span class="slide-type">{slide_type.upper()}</span>
            </div>
            <div class="slide-content">
                {''.join(body_parts)}
            </div>
        </div>
        """

    # ── PowerPoint Export ───────────────────────────────────────────────────

    def export_powerpoint(self, output_path: str = None) -> str:
        """Export deck as a .pptx file. Requires python-pptx."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")

        deck = self.generate_full_deck()
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        for num, content in deck.items():
            self._add_pptx_slide(prs, num, content)

        out = output_path or str(self.output_dir / "BrownBioTech_PitchDeck.pptx")
        prs.save(out)
        return out

    def _add_pptx_slide(self, prs: Presentation, num: int, content: str):
        import re

        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        title_match = re.search(r'^# (.+)$', content, re.MULTILINE)
        subtitle_match = re.search(r'^### (.+)$', content, re.MULTILINE)

        title = title_match.group(1) if title_match else f"Slide {num}"
        subtitle = subtitle_match.group(1) if subtitle_match else ""

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
        tf = txBox.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

        # Subtitle
        if subtitle:
            txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.33), Inches(0.5))
            tf2 = txBox2.text_frame
            tf2.paragraphs[0].text = subtitle
            tf2.paragraphs[0].font.size = Pt(18)
            tf2.paragraphs[0].font.color.rgb = RGBColor(0x4A, 0x4A, 0x6A)

        # Content bullets
        lines = [l.strip() for l in content.split("\n") if l.strip()]
        bullet_lines = [l[2:] for l in lines if l.startswith("• ") or l.startswith("- ")]
        other_lines = [l for l in lines if not l.startswith("• ") and not l.startswith("- ") and not l.startswith("#") and not l.startswith("---") and not l.startswith(">") and "## " not in l and "### " not in l]

        if bullet_lines:
            y = 1.8 if subtitle else 1.3
            txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12.33), Inches(5.0))
            tf3 = txBox3.text_frame
            tf3.word_wrap = True
            for i, bullet in enumerate(bullet_lines):
                p = tf3.paragraphs[i] if i == 0 else tf3.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(0x2D, 0x2D, 0x44)
                p.space_after = Pt(8)


# ─── HTML Template ───────────────────────────────────────────────────────────

HTML_TEMPLATE = r"""
<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>BrownBioTech — Investor Pitch Deck</title>
<style>
  @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');

  :root {
    --brown: #4A2010;
    --brown-light: #7A4020;
    --accent: #C4773A;
    --dark: #0D0D1A;
    --text: #1A1A2E;
    --muted: #6B6B8A;
    --bg: #FAFAF8;
    --card: #FFFFFF;
    --highlight: #E8F0FF;
    --border: #E0DDD8;
  }

  * { box-sizing: border-box; margin: 0; padding: 0; }

  body {
    font-family: 'Inter', -apple-system, sans-serif;
    background: var(--dark);
    color: var(--text);
    font-size: 16px;
    line-height: 1.6;
  }

  /* ── Deck wrapper ── */
  .deck {
    display: flex;
    overflow-x: auto;
    scroll-snap-type: x mandatory;
    height: 100vh;
    scroll-behavior: smooth;
  }

  /* ── Individual slide ── */
  .slide {
    min-width: 100vw;
    height: 100vh;
    scroll-snap-align: start;
    display: flex;
    flex-direction: column;
    background: var(--bg);
    position: relative;
    padding: 60px 80px;
  }

  /* ── Slide header ── */
  .slide-header {
    display: flex;
    justify-content: space-between;
    align-items: center;
    margin-bottom: 40px;
    padding-bottom: 16px;
    border-bottom: 2px solid var(--border);
  }

  .slide-num {
    font-size: 13px;
    font-weight: 600;
    color: var(--accent);
    letter-spacing: 0.1em;
    text-transform: uppercase;
  }

  .slide-type {
    font-size: 11px;
    font-weight: 600;
    color: var(--muted);
    letter-spacing: 0.15em;
    text-transform: uppercase;
  }

  /* ── Slide content ── */
  .slide-content {
    flex: 1;
    display: flex;
    flex-direction: column;
    justify-content: flex-start;
    max-width: 1100px;
  }

  .slide-content h1 {
    font-size: 52px;
    font-weight: 700;
    color: var(--brown);
    line-height: 1.1;
    margin-bottom: 12px;
  }

  .slide-content h2 {
    font-size: 38px;
    font-weight: 700;
    color: var(--brown);
    margin-bottom: 28px;
    padding-bottom: 12px;
    border-bottom: 3px solid var(--accent);
    display: inline-block;
  }

  .slide-content h3 {
    font-size: 22px;
    font-weight: 600;
    color: var(--brown-light);
    margin-bottom: 20px;
    margin-top: 8px;
  }

  .slide-content h4 {
    font-size: 18px;
    font-weight: 600;
    color: var(--text);
    margin: 16px 0 8px;
  }

  .slide-content p {
    font-size: 18px;
    line-height: 1.7;
    color: var(--text);
    margin-bottom: 14px;
  }

  .slide-content ul.bullets {
    list-style: none;
    margin: 12px 0 20px;
  }

  .slide-content ul.bullets li {
    font-size: 17px;
    line-height: 1.65;
    color: var(--text);
    padding: 8px 0 8px 28px;
    position: relative;
    border-bottom: 1px solid var(--border);
  }

  .slide-content ul.bullets li::before {
    content: '▸';
    position: absolute;
    left: 0;
    color: var(--accent);
    font-weight: 700;
  }

  .slide-content blockquote {
    background: var(--highlight);
    border-left: 4px solid var(--accent);
    padding: 16px 20px;
    border-radius: 0 8px 8px 0;
    font-size: 16px;
    color: var(--text);
    font-style: italic;
    margin: 16px 0;
  }

  .slide-content .table-row {
    font-size: 15px;
    color: var(--muted);
    font-family: monospace;
    margin: 2px 0;
  }

  /* ── Title slide ── */
  .slide-title {
    background: linear-gradient(135deg, var(--brown) 0%, #2A0F05 100%);
    color: white;
    justify-content: center;
  }

  .slide-title .slide-content h1 {
    color: #F5E6D3;
    font-size: 72px;
    margin-bottom: 8px;
  }

  .slide-title .slide-content h3 {
    color: var(--accent);
    font-size: 26px;
    font-weight: 400;
    margin-bottom: 4px;
  }

  .slide-title .slide-content p {
    color: rgba(255,255,255,0.7);
    font-size: 20px;
    margin-bottom: 40px;
  }

  .slide-title .slide-header {
    border-bottom-color: rgba(255,255,255,0.15);
  }

  .slide-title .slide-num { color: rgba(255,255,255,0.4); }
  .slide-title .slide-type { color: rgba(255,255,255,0.3); }

  .title-divider {
    width: 80px;
    height: 3px;
    background: var(--accent);
    margin: 20px 0;
  }

  .series-badge {
    display: inline-block;
    background: var(--accent);
    color: white;
    font-size: 15px;
    font-weight: 700;
    letter-spacing: 0.05em;
    padding: 8px 20px;
    border-radius: 4px;
    margin-top: 20px;
  }

  /* ── Problem slide ── */
  .slide-problem .slide-content h2::before {
    content: '⚠ ';
  }

  /* ── Solution / highlight ── */
  .slide-solution .highlight-box {
    background: linear-gradient(135deg, var(--brown), var(--brown-light));
    color: white;
    padding: 16px 24px;
    border-radius: 8px;
    font-weight: 600;
    font-size: 15px;
    margin-top: 24px;
    display: inline-block;
    letter-spacing: 0.05em;
  }

  /* ── Pipeline slide ── */
  .pipeline-table {
    width: 100%;
    border-collapse: collapse;
    margin-top: 20px;
    font-size: 15px;
  }

  .pipeline-table th {
    background: var(--brown);
    color: white;
    padding: 12px 16px;
    text-align: left;
    font-weight: 600;
    font-size: 13px;
    letter-spacing: 0.05em;
    text-transform: uppercase;
  }

  .pipeline-table td {
    padding: 14px 16px;
    border-bottom: 1px solid var(--border);
    vertical-align: top;
  }

  .pipeline-table tr:last-child td { border-bottom: none; }

  .lead-badge {
    background: var(--accent);
    color: white;
    font-size: 11px;
    font-weight: 700;
    padding: 3px 8px;
    border-radius: 3px;
    margin-left: 8px;
    vertical-align: middle;
    text-transform: uppercase;
    letter-spacing: 0.05em;
  }

  /* ── Market slide ── */
  .market-grid {
    display: grid;
    grid-template-columns: repeat(2, 1fr);
    gap: 20px;
    margin: 20px 0;
  }

  .market-card {
    background: var(--card);
    border: 1px solid var(--border);
    border-radius: 10px;
    padding: 24px;
    box-shadow: 0 2px 8px rgba(0,0,0,0.05);
  }

  .market-card .value {
    font-size: 32px;
    font-weight: 700;
    color: var(--brown);
  }

  .market-card .label {
    font-size: 14px;
    color: var(--muted);
    margin-bottom: 4px;
    text-transform: uppercase;
    letter-spacing: 0.05em;
  }

  /* ── Team slide ── */
  .founder-card {
    background: linear-gradient(135deg, var(--brown), var(--brown-light));
    color: white;
    border-radius: 12px;
    padding: 28px 32px;
    margin-bottom: 24px;
  }

  .founder-card .name {
    font-size: 22px;
    font-weight: 700;
    margin-bottom: 4px;
  }

  .founder-card .title {
    font-size: 14px;
    opacity: 0.8;
    margin-bottom: 12px;
  }

  .founder-card .detail {
    font-size: 14px;
    opacity: 0.9;
    line-height: 1.6;
  }

  /* ── Ask slide ── */
  .funds-table {
    width: 100%;
    border-collapse: collapse;
    margin: 20px 0;
  }

  .funds-table th {
    background: var(--brown);
    color: white;
    padding: 12px 16px;
    text-align: left;
    font-size: 13px;
    text-transform: uppercase;
    letter-spacing: 0.05em;
  }

  .funds-table td {
    padding: 14px 16px;
    border-bottom: 1px solid var(--border);
    font-size: 15px;
  }

  .funds-table td:nth-child(2) {
    font-weight: 700;
    color: var(--brown);
  }

  .funds-table tr:nth-child(odd) td {
    background: rgba(74, 32, 16, 0.03);
  }

  .milestone-list {
    list-style: none;
    margin-top: 20px;
  }

  .milestone-list li {
    font-size: 16px;
    padding: 10px 0;
    border-bottom: 1px solid var(--border);
    color: var(--text);
  }

  .milestone-list li::before {
    content: '→ ';
    color: var(--accent);
    font-weight: 700;
  }

  /* ── Vision quote ── */
  .vision {
    background: var(--highlight);
    border-left: 4px solid var(--brown);
    padding: 16px 20px;
    border-radius: 0 8px 8px 0;
    font-size: 16px;
    font-style: italic;
    color: var(--text);
    margin-top: 24px;
  }

  /* ── Navigation ── */
  .deck-nav {
    position: fixed;
    bottom: 24px;
    right: 24px;
    display: flex;
    gap: 8px;
    z-index: 100;
  }

  .deck-nav a {
    display: inline-flex;
    align-items: center;
    justify-content: center;
    width: 36px;
    height: 36px;
    background: var(--brown);
    color: white;
    border-radius: 50%;
    text-decoration: none;
    font-size: 18px;
    font-weight: 700;
    opacity: 0.85;
    transition: opacity 0.2s;
  }

  .deck-nav a:hover { opacity: 1; }

  /* ── Keyboard hint ── */
  .kb-hint {
    position: fixed;
    bottom: 24px;
    left: 24px;
    font-size: 12px;
    color: rgba(255,255,255,0.35);
    z-index: 100;
  }

  .slide:not(.slide-title) .kb-hint { display: none; }

  /* ── Tech highlight ── */
  .tech-badge {
    display: inline-block;
    background: var(--highlight);
    border: 1px solid rgba(74, 32, 16, 0.2);
    color: var(--brown);
    font-size: 13px;
    font-weight: 600;
    padding: 8px 16px;
    border-radius: 20px;
    margin-top: 16px;
  }

  /* ── Responsive ── */
  @media (max-width: 900px) {
    .slide { padding: 40px 32px; }
    .slide-content h1 { font-size: 40px; }
    .slide-content h2 { font-size: 28px; }
    .market-grid { grid-template-columns: 1fr; }
  }
</style>
</head>
<body>

<div class="deck" id="deck">
{slides}
</div>

<div class="deck-nav">
  <a href="javascript:void(0)" onclick="deck.scrollBy({{-width: window.innerWidth}, behavior:'smooth'}})">‹</a>
  <a href="javascript:void(0)" onclick="deck.scrollBy({{width: window.innerWidth}, behavior:'smooth'}})">›</a>
</div>

<div class="kb-hint">← → Arrow keys to navigate</div>

<script>
  // Arrow key navigation
  document.addEventListener('keydown', function(e) {{
    const deck = document.getElementById('deck');
    if (e.key === 'ArrowRight' || e.key === 'ArrowDown') {{
      deck.scrollBy({{left: window.innerWidth, behavior: 'smooth'}});
    }} else if (e.key === 'ArrowLeft' || e.key === 'ArrowUp') {{
      deck.scrollBy({{left: -window.innerWidth, behavior: 'smooth'}});
    }}
  }});

  // Scroll spy — highlight current slide
  const deck = document.getElementById('deck');
  const observer = new IntersectionObserver((entries) => {{
    entries.forEach(entry => {{
      if (entry.isIntersecting) {{
        // could track active slide here
      }}
    }});
  }}, {{ threshold: 0.5 }});

  deck.querySelectorAll('.slide').forEach(s => observer.observe(s));
</script>

</body>
</html>
"""


# ─── CLI ─────────────────────────────────────────────────────────────────────

def main():
    gen = PitchDeckGenerator()

    print("BrownBioTech Pitch Deck Generator")
    print("=" * 40)

    # Generate all slides as markdown
    deck = gen.generate_full_deck()
    for num, content in deck.items():
        out_file = gen.output_dir / f"slide_{num:02d}.md"
        out_file.write_text(content.strip(), encoding="utf-8")
        print(f"  ✓ Slide {num:02d} → {out_file.name}")

    # Export HTML
    html_path = gen.export_html()
    print(f"  ✓ HTML Deck → {Path(html_path).name}")

    # Export PowerPoint if available
    if PPTX_AVAILABLE:
        pptx_path = gen.export_powerpoint()
        print(f"  ✓ PowerPoint → {Path(pptx_path).name}")
    else:
        print("  ⚠ python-pptx not available — PowerPoint export